
from pptx.enum.chart import XL_LEGEND_POSITION


def use_line():
    from pptx import Presentation
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE,XL_DATA_LABEL_POSITION
    from pptx.util import Inches,Pt
    ppt = Presentation()
    slide = ppt.slides.add_slide(ppt.slide_layouts[6])
    # 封装图表数据
    chart_data = CategoryChartData()
    # 分组数据
    chart_data.categories = ['第一季度','第二季度','第三季度','第四季度']
    # 具体数据
    chart_data.add_series('series1',(19,21,16))
    chart_data.add_series('series2',(22,28,15))
    chart_data.add_series('series3',(20,26,18))


    x = y= Inches(2)
    width = Inches(6)
    height = Inches(4.5)
    # 绘制图表
    chart = slide.shapes.add_chart(XL_CHART_TYPE.LINE,x,y,width,height,chart_data).chart
    chart.chart_style = 10  # 1-48 主题颜色
    chart.font.size = Pt(10)

    category_axis = chart.category_axis
    category_axis.tick_labels.font.size = Pt(15)    # 设置axis的字体大小
    category_axis.has_major_gridlines =True # 是否显示为表格

    # 设置标签
    plot = chart.plots[0]
    plot.has_data_labels = True
    data_labels = plot.data_labels
    data_labels.position = XL_DATA_LABEL_POSITION.INSIDE_END

    # 设置图例
    chart.has_legend = True
    chart.legend.font.size = Pt(15)
    chart.legend.position = XL_LEGEND_POSITION.TOP  # 图例位置
    chart.legend.include_in_layout = False

    ppt.save('./create_data/31_绘制图表_折线图.pptx')
def use_pie():
    from pptx import Presentation
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE,XL_DATA_LABEL_POSITION
    from pptx.util import Inches,Pt
    ppt = Presentation()
    slide = ppt.slides.add_slide(ppt.slide_layouts[6])
    # 封装图表数据
    chart_data = CategoryChartData()
    chart_data.categories =['第一季度','第二季度','第三季度','第四季度']
    chart_data.add_series('季度销售比例',(0.27,0.23,0.31,0.19))

    x = y= Inches(2)
    width = Inches(6)
    height = Inches(4.5)
    chart = slide.shapes.add_chart(XL_CHART_TYPE.PIE,x,y,width,height,chart_data).chart
    chart.has_legend = True
    chart.plots[0].has_data_labels = True
    data_labels = chart.plots[0].data_labels
    data_labels.position = XL_DATA_LABEL_POSITION.OUTSIDE_END
    data_labels.number_format ='0%'
    
    ppt.save('./create_data/31_绘制图表_饼图.pptx')
if __name__ == "__main__":
    # use_line()
    use_pie()