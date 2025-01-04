def use_bar():
    from pptx import Presentation
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.util import Inches
    ppt = Presentation()
    slide = ppt.slides.add_slide(ppt.slide_layouts[6])
    # 封装图表数据
    chart_data = CategoryChartData()
    # 分组数据
    chart_data.categories = ['第一季度','第二季度','第三季度','第四季度']
    # 具体数据
    chart_data.add_series('series',(19,21,16))

    x = y= Inches(2)
    width = Inches(6)
    height = Inches(4.5)
    # 绘制图表
    slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,x,y,width,height,chart_data)

    ppt.save('./create_data/29_绘制图表_条图.pptx')

def use_bar2():
    from pptx import Presentation
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.util import Inches
    ppt = Presentation()
    slide = ppt.slides.add_slide(ppt.slide_layouts[6])
    # 封装图表数据
    chart_data = CategoryChartData()
    # 分组数据
    chart_data.categories = ['第一季度','第二季度','第三季度','第四季度']
    # 具体数据
    chart_data.add_series('series',(19,21,16))
    chart_data.add_series('series',(22,28,15))
    chart_data.add_series('series',(20,26,18))


    x = y= Inches(2)
    width = Inches(6)
    height = Inches(4.5)
    # 绘制图表
    slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,x,y,width,height,chart_data)

    ppt.save('./create_data/29_绘制图表_条图2.pptx')

if __name__ == "__main__":
    # use_bar()
    use_bar2()