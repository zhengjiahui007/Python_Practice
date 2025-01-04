from sys import hexversion


def base_use():
    from pptx import Presentation
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
    from pptx.util import Inches,Pt
    ppt = Presentation()
    slide = ppt.slides.add_slide(ppt.slide_layouts[5])
    shapes = slide.shapes
    shapes.title.text ='流程图'

    top = Inches(3)
    left = Inches(1)
    width = Inches(2)
    height = Inches(1)
    t_sh = shapes.add_shape(MSO_AUTO_SHAPE_TYPE.PENTAGON,left,top,width,height)
    t_sh.text = '第1步'

    for i in range(2,6):
        left = left + width - Inches(0.4)
        temp = shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON,left,top,width,height)
        # temp.text = f'第{i}步'
        p = temp.text_frame
        p.text =  f'第{i}步'
        p.fit_text(max_size =10,bold = True,italic = True)  # 必须设置内容之后
    
    ppt.save('./create_data/28_增加流程图.pptx')

if __name__ == "__main__":
    base_use()