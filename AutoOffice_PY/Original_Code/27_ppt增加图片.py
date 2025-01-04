def add_pic():
    from pptx import Presentation
    from pptx.util import Pt
    ppt = Presentation()
    slide = ppt.slides.add_slide(ppt.slide_layouts[1])
    # 获取内容框
    shapes = slide.shapes
    left = top =Pt(30)
    shapes.add_picture('./base_data/backg.jpg',left,top)

    slide2 = ppt.slides.add_slide(ppt.slide_layouts[1])
    # 获取内容框
    shapes = slide2.shapes
    left = top =Pt(0)
    height = Pt(300)
    shapes.add_picture('./base_data/backg.jpg',left,top,height)
    ppt.save('./create_data/27_ppt增加图片.pptx')

if __name__ == "__main__":
    add_pic()