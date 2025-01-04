# pip install python-pptx
def base_ppt1():
    from pptx import Presentation
    # 创建PPT
    ppt = Presentation()
    slide = ppt.slides.add_slide(ppt.slide_layouts[0])
    slide = ppt.slides.add_slide(ppt.slide_layouts[1])
    slide = ppt.slides.add_slide(ppt.slide_layouts[3])
    slide = ppt.slides.add_slide(ppt.slide_layouts[6])
    ppt.save('./create_data/26_base_ppt01.pptx')

def base_ppt2():
    from pptx import Presentation
    from pptx.util import Pt
    # 创建PPT
    ppt = Presentation()
    slide = ppt.slides.add_slide(ppt.slide_layouts[1])
    # 获取内容框
    shapes = slide.shapes
    title = shapes.title
    title.text = 'This is Python Title!!!!'
    cotent = shapes.placeholders[1]
    cotent.text ='Content Info' # content.text_frame.text = 'Content Info'

    slide2 = ppt.slides.add_slide(ppt.slide_layouts[1])
    shapes = slide2.shapes
    content = shapes.placeholders[1]
    tf = content.text_frame
    p1 = tf.add_paragraph()
    p1.text = '这个是内容1'
    p1.level =1
    p2 = tf.add_paragraph()
    p2.text = '这个是内容2'
    p2.level =2
    p3 = tf.add_paragraph()
    p3.text = '这个是内容3'
    p3.level =3


    p4 = tf.add_paragraph()
    p4.text = '这个是内容3'
    p4.font.bold= True
    p4.font.size = Pt(20)

    slide3 = ppt.slides.add_slide(ppt.slide_layouts[6])
    left = top = width = height = Pt(300)
    text_box = slide3.shapes.add_textbox(left,top,width,height)
    tf  = text_box.text_frame
    tf.text = '这个文本框内容！'
    p4 = tf.add_paragraph()
    p4.text = '这个是内容3'

    ppt.save('./create_data/26_base_ppt03.pptx')
if __name__ == "__main__":
    # base_ppt1()
    base_ppt2()